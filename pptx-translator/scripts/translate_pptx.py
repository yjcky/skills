#!/usr/bin/env python3
"""
PPTX Translator - Translate PowerPoint files using Amazon Translate, Bedrock LLM,
Google Translate, or Cerebras LLM.

Usage:
    python translate_pptx.py input.pptx --source zh --target en
    python translate_pptx.py input.pptx --source en --target zh --engine bedrock
    python translate_pptx.py input.pptx --source zh --target en --engine cerebras
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Import shared translation infrastructure
from translate_common import (
    TranslationEngine, AmazonTranslateEngine, BedrockLLMEngine,
    GoogleTranslateEngine, CerebrasEngine,
    create_engine, load_glossary, import_terminology,
    post_process_translation, is_cjk_char,
    LANGUAGE_NAMES,
    safe_resolve_path, setup_windows_encoding, add_common_arguments,
    build_engine_from_args,
)

# Language code to MSO_LANGUAGE_ID mapping (pptx-specific)
LANGUAGE_CODE_TO_LANGUAGE_ID = {
    'af': MSO_LANGUAGE_ID.AFRIKAANS, 'am': MSO_LANGUAGE_ID.AMHARIC,
    'ar': MSO_LANGUAGE_ID.ARABIC, 'bg': MSO_LANGUAGE_ID.BULGARIAN,
    'bn': MSO_LANGUAGE_ID.BENGALI, 'bs': MSO_LANGUAGE_ID.BOSNIAN,
    'cs': MSO_LANGUAGE_ID.CZECH, 'da': MSO_LANGUAGE_ID.DANISH,
    'de': MSO_LANGUAGE_ID.GERMAN, 'el': MSO_LANGUAGE_ID.GREEK,
    'en': MSO_LANGUAGE_ID.ENGLISH_US, 'es': MSO_LANGUAGE_ID.SPANISH,
    'et': MSO_LANGUAGE_ID.ESTONIAN, 'fi': MSO_LANGUAGE_ID.FINNISH,
    'fr': MSO_LANGUAGE_ID.FRENCH, 'fr-CA': MSO_LANGUAGE_ID.FRENCH_CANADIAN,
    'ha': MSO_LANGUAGE_ID.HAUSA, 'he': MSO_LANGUAGE_ID.HEBREW,
    'hi': MSO_LANGUAGE_ID.HINDI, 'hr': MSO_LANGUAGE_ID.CROATIAN,
    'hu': MSO_LANGUAGE_ID.HUNGARIAN, 'id': MSO_LANGUAGE_ID.INDONESIAN,
    'it': MSO_LANGUAGE_ID.ITALIAN, 'ja': MSO_LANGUAGE_ID.JAPANESE,
    'ka': MSO_LANGUAGE_ID.GEORGIAN, 'ko': MSO_LANGUAGE_ID.KOREAN,
    'lv': MSO_LANGUAGE_ID.LATVIAN, 'ms': MSO_LANGUAGE_ID.MALAYSIAN,
    'nl': MSO_LANGUAGE_ID.DUTCH, 'no': MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
    'pl': MSO_LANGUAGE_ID.POLISH, 'ps': MSO_LANGUAGE_ID.PASHTO,
    'pt': MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE, 'ro': MSO_LANGUAGE_ID.ROMANIAN,
    'ru': MSO_LANGUAGE_ID.RUSSIAN, 'sk': MSO_LANGUAGE_ID.SLOVAK,
    'sl': MSO_LANGUAGE_ID.SLOVENIAN, 'so': MSO_LANGUAGE_ID.SOMALI,
    'sq': MSO_LANGUAGE_ID.ALBANIAN, 'sr': MSO_LANGUAGE_ID.SERBIAN_LATIN,
    'sv': MSO_LANGUAGE_ID.SWEDISH, 'sw': MSO_LANGUAGE_ID.SWAHILI,
    'ta': MSO_LANGUAGE_ID.TAMIL, 'th': MSO_LANGUAGE_ID.THAI,
    'tr': MSO_LANGUAGE_ID.TURKISH, 'uk': MSO_LANGUAGE_ID.UKRAINIAN,
    'ur': MSO_LANGUAGE_ID.URDU, 'vi': MSO_LANGUAGE_ID.VIETNAMESE,
    'zh': MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
    'zh-TW': MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
}


# ---------------------------------------------------------------------------
# PPTX-specific: font auto-fit
# ---------------------------------------------------------------------------

def estimate_text_width_emu(text, font_size_emu):
    """Estimate rendered line width in EMU using character-width heuristics.

    CJK/fullwidth chars are ~2× the width of Latin chars at the same point size.
    """
    width = 0
    for c in text:
        if is_cjk_char(c):
            width += font_size_emu * 1.0
        elif c in (' ', '\t'):
            width += font_size_emu * 0.15
        else:
            width += font_size_emu * 0.5
    return int(width)


def _get_available_text_width(shape):
    """Calculate available text width in EMU, accounting for shape margins."""
    tf = shape.text_frame
    margin_left = tf.margin_left if tf.margin_left is not None else 91440
    margin_right = tf.margin_right if tf.margin_right is not None else 91440
    return max(0, shape.width - margin_left - margin_right)


def auto_fit_paragraph_font(shape, para, min_pt=8, original_text=None):
    """Shrink font size so translated text fits within the shape width.

    For paragraphs with explicit font sizes, compares estimated text width
    against available shape width directly. For inherited sizes, uses a
    font-size-independent comparison between original CJK and translated EN
    text widths to determine the shrink ratio. Preserves proportional size
    relationships between runs.
    """
    if not shape or not para:
        return
    text = para.text
    if not text or not text.strip():
        return

    from pptx.util import Pt

    ref_size_emu = None
    for run in para.runs:
        if run.font.size is not None:
            ref_size_emu = run.font.size
            break
    if ref_size_emu is None:
        try:
            if para.font.size is not None:
                ref_size_emu = para.font.size
        except Exception:
            pass
    if ref_size_emu is None:
        try:
            from pptx.oxml.ns import qn
            pPr = para._p.find(qn('a:pPr'))
            if pPr is not None:
                defRPr = pPr.find(qn('a:defRPr'))
                if defRPr is not None:
                    sz = defRPr.get('sz')
                    if sz:
                        ref_size_emu = int(int(sz) * 127)
        except Exception:
            pass

    available = _get_available_text_width(shape)
    if available <= 0:
        return

    if ref_size_emu and ref_size_emu > 0:
        # Explicit font size known — use absolute width comparison
        current_pt = ref_size_emu / 12700
        lines = text.split('\n')
        max_line_width = 0
        for line in lines:
            if line.strip():
                w = estimate_text_width_emu(line, ref_size_emu)
                if w > max_line_width:
                    max_line_width = w
        allowed = available * 1.3  # 30% overflow tolerance
        if max_line_width <= allowed:
            return
        ratio = allowed / max_line_width
    elif original_text and original_text.strip():
        # Inherited font size — use font-size-independent CJK/EN comparison
        ref_emu = 12700  # 1pt as neutral reference
        cjk_w = estimate_text_width_emu(original_text, ref_emu)
        en_w = estimate_text_width_emu(text, ref_emu)
        if en_w <= cjk_w * 1.3:
            return
        ratio = cjk_w / en_w
        ratio = ratio ** 0.5  # sqrt-soften: wrapping helps absorb extra width
        shape_h_pt = shape.height / 12700
        current_pt = max(10, shape_h_pt * 0.25)
    else:
        return

    for run in para.runs:
        if run.font.size is not None:
            run_pt = run.font.size / 12700
            run.font.size = int(max(min_pt, run_pt * ratio) * 12700)
        else:
            run.font.size = int(max(min_pt, current_pt * ratio) * 12700)


# ---------------------------------------------------------------------------
# PPTX-specific: text extraction & application
# ---------------------------------------------------------------------------

def iter_shapes(shapes):
    """Recursively traverse shapes, expanding GroupShapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def extract_texts_from_presentation(presentation) -> list:
    """Extract all translatable texts with their locations."""
    texts = []

    for slide_idx, slide in enumerate(presentation.slides):
        for shape_idx, shape in enumerate(iter_shapes(slide.shapes)):
            if shape.has_table:
                for row_idx, row in enumerate(shape.table.rows):
                    for cell_idx, cell in enumerate(row.cells):
                        for para_idx, para in enumerate(cell.text_frame.paragraphs):
                            para_text = ''.join(run.text for run in para.runs if run.text.strip())
                            if para_text.strip():
                                texts.append({
                                    'text': para_text,
                                    'location': ('table', slide_idx, shape_idx, row_idx, cell_idx, para_idx, 'paragraph'),
                                    'original_text': para_text,
                                    'paragraph': para,
                                    'shape': shape
                                })
            elif shape.has_text_frame:
                for para_idx, para in enumerate(shape.text_frame.paragraphs):
                    para_text = ' '.join(run.text.strip() for run in para.runs if run.text.strip())
                    if para_text.strip():
                        texts.append({
                            'text': para_text,
                            'location': ('shape', slide_idx, shape_idx, para_idx, 'paragraph'),
                            'original_text': para_text,
                            'paragraph': para,
                            'shape': shape
                        })

        if slide.has_notes_slide:
            for para_idx, para in enumerate(slide.notes_slide.notes_text_frame.paragraphs):
                para_text = ' '.join(run.text.strip() for run in para.runs if run.text.strip())
                if para_text.strip():
                    texts.append({
                        'text': para_text,
                        'location': ('notes', slide_idx, para_idx, 'paragraph'),
                        'original_text': para_text,
                        'paragraph': para,
                        'shape': None
                    })

    return texts


def find_run_by_location(location, presentation):
    """Find the specific run object based on location information."""
    loc_type, slide_idx, *rest = location

    slide = presentation.slides[slide_idx]
    shapes = list(iter_shapes(slide.shapes))

    if loc_type == 'shape':
        shape_idx, para_idx, run_type = rest
        if shape_idx >= len(shapes):
            return None
        shape = shapes[shape_idx]
        if shape.has_text_frame:
            para = shape.text_frame.paragraphs[para_idx]
            if run_type == 'paragraph':
                return para.runs[0] if para.runs else None
    elif loc_type == 'table':
        shape_idx, row_idx, cell_idx, para_idx, run_idx = rest
        if shape_idx >= len(shapes):
            return None
        shape = shapes[shape_idx]
        if shape.has_table:
            row = shape.table.rows[row_idx]
            cell = row.cells[cell_idx]
            para = cell.text_frame.paragraphs[para_idx]
            if run_idx < len(para.runs):
                return para.runs[run_idx]
    elif loc_type == 'notes':
        para_idx, run_type = rest
        if slide.has_notes_slide:
            para = slide.notes_slide.notes_text_frame.paragraphs[para_idx]
            if run_type == 'paragraph':
                return para.runs[0] if para.runs else None

    return None


def apply_paragraph_translation(para, translated_text):
    """Apply translated text to a paragraph — first run gets all text, rest cleared."""
    if not para or not translated_text:
        return

    runs = para.runs
    if runs:
        runs[0].text = translated_text
        for run in runs[1:]:
            run.text = ""
    else:
        para.add_run(translated_text)


# ---------------------------------------------------------------------------
# Main translation logic
# ---------------------------------------------------------------------------

def translate_presentation(presentation, engine: TranslationEngine, source_lang: str,
                          target_lang: str, batch_mode: bool = False):
    """Translate all text in a presentation."""

    if batch_mode and (isinstance(engine, (BedrockLLMEngine, CerebrasEngine))):
        print("Extracting texts...")
        text_items = extract_texts_from_presentation(presentation)

        if not text_items:
            print("No text found to translate.")
            return

        print(f"Found {len(text_items)} text segments")

        texts = [item['text'] for item in text_items]
        batch_size = engine.batch_size

        batches = []
        for i in range(0, len(texts), batch_size):
            batches.append(texts[i:i + batch_size])
        batch_count = len(batches)

        import concurrent.futures
        translated_texts = []

        max_workers = min(3, batch_count)
        print(f"Translating {batch_count} batches with {max_workers} concurrent workers...")

        results_by_index = {}
        failed_count = 0

        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_idx = {}
            for idx, batch in enumerate(batches):
                future = executor.submit(engine.translate_batch, batch, source_lang, target_lang)
                future_to_idx[future] = idx

            for future in concurrent.futures.as_completed(future_to_idx):
                idx = future_to_idx[future]
                try:
                    results_by_index[idx] = future.result()
                except Exception as e:
                    failed_count += 1
                    print(f"  [ERROR] Batch {idx + 1}/{batch_count} failed: {e}")
                    results_by_index[idx] = [str(t) for t in batches[idx]]

            for i in range(batch_count):
                batch_result = results_by_index.get(i, [str(t) for t in batches[i]])
                translated_texts.extend(batch_result)

        if batch_count > 0 and failed_count == batch_count:
            raise RuntimeError(
                f"All {batch_count} translation batches failed. "
                "Check network connectivity to Cerebras API."
            )

        print("Applying translations...")
        if len(text_items) != len(translated_texts):
            print(f"  [WARN] Translation count mismatch: {len(text_items)} original vs {len(translated_texts)} translated")
            while len(translated_texts) < len(text_items):
                original_item = text_items[len(translated_texts)]
                translated_texts.append(str(original_item['text']) if original_item['text'] else "")

        for idx, (item, translated) in enumerate(zip(text_items, translated_texts)):
            if not translated or len(str(translated).strip()) == 0:
                print(f"  [WARN] Empty translation at position {idx}, using original text")
                translated = str(item['text']) if item['text'] else ""

            if isinstance(translated, dict):
                print(f"  [WARN] Dictionary format at position {idx}, converting to string")
                translated_text = str(translated)
            elif isinstance(translated, list):
                print(f"  [WARN] List format at position {idx}, joining to string")
                translated_text = ' '.join(str(t) for t in translated)
            else:
                translated_text = str(translated)

            translated_text = post_process_translation(translated_text, target_lang)

            if 'paragraph' in item:
                para = item['paragraph']
                try:
                    apply_paragraph_translation(para, translated_text)
                    if item.get('shape'):
                        try:
                            auto_fit_paragraph_font(item['shape'], para, original_text=item.get('original_text'))
                        except Exception:
                            pass
                    print(f"  Applied paragraph translation to {item['location'][0]} on slide {item['location'][1] + 1}")
                except Exception as e:
                    print(f"  [ERROR] Failed to apply paragraph translation at position {idx}: {e}")
                    apply_paragraph_translation(para, item['original_text'] if item['original_text'] else "")
            else:
                run = find_run_by_location(item['location'], presentation)
                if run:
                    try:
                        run.text = translated_text
                        print(f"  Applied translation to {item['location'][0]} on slide {item['location'][1] + 1}")
                    except Exception as e:
                        print(f"  [ERROR] Failed to apply translation at position {idx}: {e}")
                        run.text = item['original_text'] if item['original_text'] else ""
                else:
                    print(f"  [ERROR] Could not find run at position {idx} with location {item['location']}")
    else:
        total_slides = len(presentation.slides)
        for slide_idx, slide in enumerate(presentation.slides, start=1):
            print(f"Slide {slide_idx}/{total_slides}")

            for shape in iter_shapes(slide.shapes):
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            translate_text_frame(cell.text_frame, engine, source_lang, target_lang, shape=shape)
                elif shape.has_text_frame:
                    translate_text_frame(shape.text_frame, engine, source_lang, target_lang, shape=shape)

            if slide.has_notes_slide:
                translate_text_frame(slide.notes_slide.notes_text_frame, engine, source_lang, target_lang)


def translate_text_frame(text_frame, engine: TranslationEngine, source_lang: str, target_lang: str, shape=None):
    """Translate all text in a text frame."""
    for paragraph in text_frame.paragraphs:
        translated = False
        for run in paragraph.runs:
            if run.text.strip():
                translated_text = engine.translate(run.text, source_lang, target_lang)
                run.text = post_process_translation(translated_text, target_lang)
                translated = True
        if translated and shape:
            try:
                auto_fit_paragraph_font(shape, paragraph)
            except Exception:
                pass


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description='Translate PPTX files using Amazon Translate, Bedrock LLM, Google Translate, or Cerebras LLM'
    )
    parser.add_argument('input_file', help='Input PPTX file path')
    add_common_arguments(parser)

    args = parser.parse_args()

    setup_windows_encoding()

    input_path = safe_resolve_path(args.input_file)
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input_file}")
        sys.exit(1)

    if args.output:
        output_path = safe_resolve_path(args.output).parent / Path(args.output).name
        output_path.parent.mkdir(parents=True, exist_ok=True)
    else:
        output_path = input_path.with_name(f"{input_path.stem}-{args.target}.pptx")

    engine, batch_mode = build_engine_from_args(args)

    print(f"Loading {input_path}...")
    presentation = Presentation(str(input_path))

    print(f"Translating from {args.source} to {args.target}...")
    try:
        translate_presentation(presentation, engine, args.source, args.target, batch_mode=batch_mode)
    except RuntimeError as e:
        print(f"Error: {e}")
        sys.exit(1)

    print(f"Saving {output_path}...")
    presentation.save(str(output_path.resolve()))
    print("Done!")


if __name__ == '__main__':
    main()
